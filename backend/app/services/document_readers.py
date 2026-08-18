"""
Reading documents: Word, Excel, PowerPoint, PDF, CSV and plain text.

Every reader returns the same shape — plain text for search and AI, plus an
``outline`` describing the structure (sheets, headings, pages) so Cerebro can
say "the Q3 tab, row 14" rather than only quoting a wall of text.

The parsing libraries are optional dependencies. A missing one produces a clear
"install this" message rather than an import error at boot, matching how the AI
and Qdrant extras behave.
"""

import csv
import io
import json
import re
from pathlib import Path
from typing import Any, Dict, List, Optional

from app.core.config import settings

TEXT_LIMIT = 200_000
PREVIEW_LIMIT = 8_000

EXTENSIONS = {
    ".docx": "docx", ".docm": "docx", ".dotx": "docx",
    ".xlsx": "xlsx", ".xlsm": "xlsx", ".xltx": "xlsx",
    ".pptx": "pptx", ".potx": "pptx",
    ".pdf": "pdf",
    ".csv": "csv", ".tsv": "csv",
    ".txt": "text", ".md": "text", ".log": "text", ".json": "text", ".xml": "text",
}

#: Old binary formats need Office itself or a converter; say so plainly rather
#: than failing with a confusing zip error.
LEGACY_EXTENSIONS = {
    ".doc": "Word 97-2003", ".xls": "Excel 97-2003", ".ppt": "PowerPoint 97-2003",
}

INSTALL_HINT = "pip install -r backend/requirements-documents.txt"


class DocumentError(RuntimeError):
    """Raised with a message intended to be shown to the user."""


def classify(path: Path) -> Optional[str]:
    return EXTENSIONS.get(path.suffix.lower())


def is_supported(path: Path) -> bool:
    return classify(path) is not None


def is_open_in_office(path: Path) -> bool:
    """
    True when Word or Excel has the file open.

    Office writes a lock file beside the original while it is open — but the
    name is not simply ``~$`` plus the filename. For a base name longer than
    eight characters Office drops the leading characters to keep the lock name
    short, so ``Case notes 2026.docx`` locks as ``~$se notes 2026.docx``.
    Checking only the exact name meant the guard almost never fired, and Cerebro
    would happily write to a document Word was about to save over.

    So: check the exact name, then any ``~$`` file whose name is a suffix of the
    original.
    """
    try:
        if (path.parent / f"~${path.name}").exists():
            return True
        for lock in path.parent.glob("~$*"):
            candidate = lock.name[2:]
            if candidate and path.name.endswith(candidate):
                return True
    except OSError:
        pass
    return False


def _check_size(path: Path) -> None:
    try:
        megabytes = path.stat().st_size / (1024 * 1024)
    except OSError as exc:
        raise DocumentError(f"Cannot read {path.name}: {exc}") from exc
    if megabytes > settings.DOCUMENT_MAX_MB:
        raise DocumentError(
            f"{path.name} is {megabytes:.1f} MB, over the "
            f"{settings.DOCUMENT_MAX_MB:.0f} MB limit. Raise it in Settings → Documents."
        )


# -------------------------------------------------------------------- Word
def read_docx(path: Path) -> Dict[str, Any]:
    try:
        import docx
    except ImportError as exc:
        raise DocumentError(f"Reading Word files needs python-docx. {INSTALL_HINT}") from exc

    document = docx.Document(str(path))

    lines: List[str] = []
    headings: List[Dict[str, Any]] = []

    for index, paragraph in enumerate(document.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name if paragraph.style else "") or ""
        if style.lower().startswith("heading") or style.lower() == "title":
            headings.append({"text": text, "style": style, "paragraph": index})
            lines.append(f"\n{text}\n")
        else:
            lines.append(text)

    tables = []
    for table_index, table in enumerate(document.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        tables.append({"index": table_index, "rows": len(rows),
                       "columns": len(rows[0]) if rows else 0,
                       "preview": rows[:5]})
        for row in rows:
            lines.append(" | ".join(row))

    return {
        "kind": "docx",
        "text": "\n".join(lines)[:TEXT_LIMIT],
        "outline": {
            "paragraphs": len(document.paragraphs),
            "headings": headings[:60],
            "tables": tables,
            "words": sum(len(line.split()) for line in lines),
        },
    }


# ------------------------------------------------------------------- Excel
def read_xlsx(path: Path, max_rows: int = 400, max_cols: int = 40) -> Dict[str, Any]:
    try:
        import openpyxl
    except ImportError as exc:
        raise DocumentError(f"Reading Excel files needs openpyxl. {INSTALL_HINT}") from exc

    # data_only reads the last-calculated values rather than formula strings —
    # what the user sees on screen, which is what questions are asked about.
    workbook = openpyxl.load_workbook(str(path), data_only=True, read_only=True)

    lines: List[str] = []
    sheets: List[Dict[str, Any]] = []

    try:
        for worksheet in workbook.worksheets:
            rows_read = 0
            headers: List[str] = []
            lines.append(f"\n## Sheet: {worksheet.title}\n")

            for row in worksheet.iter_rows(max_row=max_rows, max_col=max_cols,
                                           values_only=True):
                values = ["" if value is None else str(value) for value in row]
                while values and not values[-1]:
                    values.pop()
                if not values:
                    continue
                if not headers:
                    headers = values
                lines.append(" | ".join(values))
                rows_read += 1

            sheets.append({
                "name": worksheet.title,
                "rows": worksheet.max_row,
                "columns": worksheet.max_column,
                "rows_read": rows_read,
                "headers": headers[:max_cols],
                "truncated": bool(worksheet.max_row and worksheet.max_row > max_rows),
            })
    finally:
        workbook.close()

    return {
        "kind": "xlsx",
        "text": "\n".join(lines)[:TEXT_LIMIT],
        "outline": {"sheets": sheets, "sheet_count": len(sheets)},
    }


# -------------------------------------------------------------- PowerPoint
def read_pptx(path: Path) -> Dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentError(
            f"Reading PowerPoint files needs python-pptx. {INSTALL_HINT}") from exc

    presentation = Presentation(str(path))
    lines: List[str] = []
    slides: List[Dict[str, Any]] = []

    for number, slide in enumerate(presentation.slides, start=1):
        texts = [shape.text.strip() for shape in slide.shapes
                 if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        title = texts[0] if texts else f"Slide {number}"
        slides.append({"number": number, "title": title[:120], "blocks": len(texts)})
        lines.append(f"\n## Slide {number}: {title}\n")
        lines.extend(texts[1:])

    return {
        "kind": "pptx",
        "text": "\n".join(lines)[:TEXT_LIMIT],
        "outline": {"slides": slides, "slide_count": len(slides)},
    }


# --------------------------------------------------------------------- PDF
def read_pdf(path: Path) -> Dict[str, Any]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentError(f"Reading PDFs needs pypdf. {INSTALL_HINT}") from exc

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        # decrypt() reports failure with a falsy return code rather than raising,
        # so an unchecked call leaves an unreadable reader and the error surfaces
        # later as the misleading "probably a scan".
        try:
            unlocked = reader.decrypt("")
        except Exception:
            unlocked = 0
        if not unlocked:
            raise DocumentError(
                f"{path.name} is password-protected. Cerebro cannot open it.")

    lines: List[str] = []
    pages: List[Dict[str, Any]] = []

    for number, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        pages.append({"number": number, "characters": len(text)})
        if text:
            lines.append(f"\n## Page {number}\n{text}")

    body = "\n".join(lines)
    if not body.strip():
        raise DocumentError(
            f"{path.name} has no extractable text — it is probably a scan. "
            "Cerebro cannot OCR it yet."
        )

    return {
        "kind": "pdf",
        "text": body[:TEXT_LIMIT],
        "outline": {"pages": pages, "page_count": len(reader.pages)},
    }


# --------------------------------------------------------- CSV / plain text
def read_csv(path: Path, max_rows: int = 500) -> Dict[str, Any]:
    raw = path.read_text(encoding="utf-8-sig", errors="replace")
    delimiter = "\t" if path.suffix.lower() == ".tsv" else None
    if delimiter is None:
        try:
            delimiter = csv.Sniffer().sniff(raw[:4000], delimiters=",;\t|").delimiter
        except csv.Error:
            delimiter = ","

    rows = list(csv.reader(io.StringIO(raw), delimiter=delimiter))
    headers = rows[0] if rows else []
    lines = [" | ".join(row) for row in rows[:max_rows]]

    return {
        "kind": "csv",
        "text": "\n".join(lines)[:TEXT_LIMIT],
        "outline": {
            "rows": len(rows), "columns": len(headers),
            "headers": headers[:60], "delimiter": delimiter,
            "truncated": len(rows) > max_rows,
        },
    }


def read_text(path: Path) -> Dict[str, Any]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    headings = [line.strip() for line in raw.splitlines()
                if line.startswith("#") or re.match(r"^[A-Z][A-Za-z ]{3,60}:$", line.strip())]
    return {
        "kind": "text",
        "text": raw[:TEXT_LIMIT],
        "outline": {"lines": raw.count("\n") + 1, "characters": len(raw),
                    "headings": headings[:40]},
    }


READERS = {
    "docx": read_docx, "xlsx": read_xlsx, "pptx": read_pptx,
    "pdf": read_pdf, "csv": read_csv, "text": read_text,
}


def read(path: Path) -> Dict[str, Any]:
    """Read any supported document into text plus an outline."""
    path = Path(path)
    if not path.exists():
        raise DocumentError(f"File not found: {path}")
    if not path.is_file():
        raise DocumentError(f"Not a file: {path}")

    legacy = LEGACY_EXTENSIONS.get(path.suffix.lower())
    if legacy:
        raise DocumentError(
            f"{path.name} is a {legacy} file. Save it as the modern format "
            f"(.docx/.xlsx/.pptx) and Cerebro can read it."
        )

    kind = classify(path)
    if kind is None:
        raise DocumentError(f"Cerebro does not read {path.suffix or 'files without an extension'} yet.")

    _check_size(path)

    result = READERS[kind](path)
    result["path"] = str(path)
    result["name"] = path.name
    result["preview"] = result["text"][:PREVIEW_LIMIT]
    result["outline_json"] = json.dumps(result["outline"], default=str)[:20_000]
    try:
        result["mtime"] = path.stat().st_mtime
        result["size_bytes"] = path.stat().st_size
    except OSError:
        result["mtime"] = None
        result["size_bytes"] = None
    return result


def available_formats() -> Dict[str, Any]:
    """Which document types can actually be read right now."""
    from importlib import util

    modules = {"docx": "docx", "xlsx": "openpyxl", "pptx": "pptx", "pdf": "pypdf"}
    formats = {}
    for kind, module in modules.items():
        try:
            formats[kind] = util.find_spec(module) is not None
        except (ImportError, ValueError):
            formats[kind] = False
    formats["csv"] = True
    formats["text"] = True
    return formats
